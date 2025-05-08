import logging
import os
from typing import List
from pptx import Presentation
import pandas as pd
from dependency_injector.wiring import inject, Provide
from langchain_community.document_loaders import (
    TextLoader,
    WebBaseLoader,
    UnstructuredWordDocumentLoader, CSVLoader, AzureAIDocumentIntelligenceLoader,
)
from langchain_community.vectorstores import Chroma
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from agents.grader_agent import GraderNode
from commands.document_qa_chat import AdvancedDocumentQAAgent
from common.aws_fs_helper import AwsS3FsHelper
from container import Container
from entities.server_entities import QueryResponse, QueryRequest, KnowledgeStore
from repositories.chroma_db_repo import DocRepository, openai_embeddings
from repositories.document_qa_repo import DocumentQADBRepository
from .ocr_service import get_ocr_text

SAVE_DIR = "downloaded_docs"
CHROMA_PERSIST_DIR = "chroma_db"


class ExcelLoader(BaseLoader):
    def __init__(self, file_path: str, sheet_name: str = None):
        self.file_path = file_path
        self.sheet_name = sheet_name

    def load(self) -> List[Document]:
        df = pd.read_excel(self.file_path, sheet_name=self.sheet_name)
        documents = []

        if isinstance(df, dict):  # Multiple sheets
            for name, sheet in df.items():
                text = sheet.to_string(index=False)
                documents.append(Document(page_content=text, metadata={"sheet_name": name, "source": self.file_path}))
        else:
            text = df.to_string(index=False)
            documents.append(Document(page_content=text, metadata={"sheet_name": self.sheet_name or 'default', "source": self.file_path}))

        return documents

class PowerPointLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        prs = Presentation(self.file_path)
        documents = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            slide_text = "\n".join(texts)
            documents.append(Document(page_content=slide_text, metadata={"slide_number": i + 1, "source": self.file_path}))

        return documents

class HandleDocumentChunks:
    def __init__(
        self, s3_helper: AwsS3FsHelper, file_path: str, doc_id: int, file_type: str
    ):
        self.s3_helper = s3_helper
        self.file_path = file_path
        self.doc_id = doc_id
        self.file_type = file_type
        self.local_file_path = None
        self.log = logging.getLogger(f"{__name__}.{self.__class__.__name__}")

    def create_local_path(self):
        return os.path.join(SAVE_DIR, os.path.basename(self.file_path))

    def get_loaded_document(self):
        self.log.info("Document loader executed.")

        if self.file_type == "url":
            return WebBaseLoader(web_path=self.file_path).load()

        ext = os.path.splitext(self.file_path)[-1]
        self.local_file_path = self.create_local_path()
        print(self.file_path, self.local_file_path)
        self.s3_helper.download_to(f"s3://{self.file_path}", self.local_file_path)
        if ext == ".docx":
            return UnstructuredWordDocumentLoader(
                file_path=self.local_file_path, mode="elements", strategy="fast"
            ).load()
        if ext == ".pdf":
            docs = Document(
                page_content=get_ocr_text(self.local_file_path),
                metadata={"source": f"s3://{self.file_path}"},
            )
            return [docs]
        if ext == ".txt":
            return TextLoader(self.local_file_path).load()

        if ext == ".csv":
            return CSVLoader(self.local_file_path).load()

        if ext == '.xlsx':
            return ExcelLoader(self.local_file_path).load()

        if ext == '.pptx':
            return PowerPointLoader(self.local_file_path).load()

    def add_metadata_document(self, documents: Document):
        self.log.info("Add Metadata Document executed.")
        processed_docs = []

        for doc in documents:
            metadata = doc.metadata
            source = metadata.get("source", "")
            doc_id = self.doc_id

            processed_docs.append(
                Document(
                    metadata={"source": source, "doc_id": doc_id},
                    page_content=doc.page_content,
                )
            )
        return processed_docs

    def split_document_chunks(self, documents):
        self.log.info("Split Document Chunk executed.")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500, chunk_overlap=350
        )
        return text_splitter.split_documents(documents)

    def get_document_chunks(self):
        documents = self.get_loaded_document()
        if self.file_type == "file":
            os.remove(self.local_file_path)
        documents = self.split_document_chunks(documents)
        return self.add_metadata_document(documents)


class DocumentService:
    @inject
    def __init__(
        self,
        s3_helper: AwsS3FsHelper = Provide[Container.s3_helper],
        doc_repo: DocRepository = Provide[Container.doc_repo],
    ):
        self.doc_repo = doc_repo
        self.s3_helper = s3_helper

    @staticmethod
    def format_file_path(file_path, file_type):
        return f"s3://{file_path}" if file_type == "file" else file_path

    def get_chroma_collection(self, index_id: str):
        return Chroma(
            persist_directory=os.path.join(CHROMA_PERSIST_DIR, index_id),
            embedding_function=openai_embeddings,
        )

    def upload_input_docs(
        self,
        index_id: str,
        index_name: str,
        operation: str,
        file_type: str,
        file_data: List,
    ):

        for file_info in file_data:
            doc_id = file_info.fileId
            doc_chunks = HandleDocumentChunks(
                self.s3_helper, file_info.filePath, doc_id, file_type
            ).get_document_chunks()

            self.doc_repo.upload_document(
                self.format_file_path(file_info.filePath, file_type),
                doc_id,
                index_id,
                index_name,
                operation,
                file_type,
                doc_chunks,
            )
        return True

    def get_all_docs(self, index_id: str):
        return self.doc_repo.get_all_docs(index_id)

    def delete_docs(self, index_id: str, file_data: List):
        for file_info in file_data:
            doc_id = file_info.fileId
            self.doc_repo.delete_by_id(doc_id, index_id)
        return True

    def delete_index(self, index_id: str):
        return self.doc_repo.delete_index(index_id)


def get_knowledge_store_format(found_doc_ids, knowledge_data) -> KnowledgeStore:
    return KnowledgeStore(
        id=knowledge_data.id,
        name=knowledge_data.name,
        type=knowledge_data.type,
        documentIds=found_doc_ids,
    )


def process_entire_document(file_path, query, source_file_path, doc_id):
    extension = file_path.split(".")[-1]
    doc_content = None

    if extension == "docx":
        loader = UnstructuredWordDocumentLoader(file_path=file_path)
        documents = loader.load()
        doc_content = "\n\n".join([doc.page_content for doc in documents])
    elif extension == "pdf":
        doc_content = get_ocr_text(file_path)
    elif extension == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            doc_content = f.read()

    qa_agent = AdvancedDocumentQAAgent()
    response = qa_agent.run_query_on_entire_document(
        query=query.question,
        metadata={"source": source_file_path, "doc_id": doc_id},
        doc_content=doc_content,
        chat_history=query.historyList,
    )
    return response


class DocChatService:
    @inject
    def __init__(
        self,
        s3_helper: AwsS3FsHelper = Provide[Container.s3_helper],
        doc_repo: DocRepository = Provide[Container.doc_repo],
        document_repo: DocumentQADBRepository = Provide[Container.document_qa_repo],
    ):
        self.doc_repo = doc_repo
        self.s3_helper = s3_helper
        self.document_repo = document_repo
        self.log = logging.getLogger(f"{__name__}.{self.__class__.__name__}")

    def document_chat(self, query: QueryRequest) -> QueryResponse:
        qa_agent = AdvancedDocumentQAAgent()
        query_answers, raw_contexts, knowledge_stores = [], [], []
        # print(query.historyList)
        for knowledge_data in query.knowledgeStoreList:
            llm_response = self.fetch_llm_response(qa_agent, query, knowledge_data)
            print(f"LLM Response: {llm_response}")
            grader = GraderNode()
            grader_response = grader(query.question, llm_response.get("output"))

            if not grader_response:
                raise ValueError("Grader response is None or empty")

            print(f"Grader Response: {grader_response}")

            response = None
            grader_response_for_whole_doc = None
            found_doc_ids = []

            if self.is_incorrect_response(grader_response):
                found_doc_ids = llm_response.get("found_doc_ids", [])
                self.log.info(f"Doc ID List: {found_doc_ids}")

                if found_doc_ids:
                    self.log.info(
                        "Fetching LLM response by processing entire document..."
                    )
                    response, grader_response_for_whole_doc, found_doc_ids = (
                        self.get_llm_response_by_processing_whole_document(
                            doc_id_list=found_doc_ids, query=query, grader=grader
                        )
                    )
            else:
                found_doc_ids = llm_response.get("found_doc_ids", [])

            final_answer = (
                grader_response_for_whole_doc.answer
                if grader_response_for_whole_doc
                else grader_response.answer
            )
            final_raw_context = (
                grader_response_for_whole_doc.raw_context
                if grader_response_for_whole_doc
                else grader_response.raw_context
            )

            query_answers.append(final_answer)
            raw_contexts.append(final_raw_context)

            self.log.info(
                f"Final Answer: {final_answer}\n"
                f"Raw Context: {final_raw_context}\n"
                f"Found Doc IDs: {found_doc_ids}"
            )

            final_doc_ids = []
            if final_answer != "No Data Found":
                final_doc_ids = found_doc_ids
            else:
                for i in found_doc_ids:
                    if i not in final_doc_ids:
                        final_doc_ids.append(i)
            self.log.info(f"Final Document IDs: {final_doc_ids}")

            knowledge_stores.append(
                get_knowledge_store_format(final_doc_ids, knowledge_data)
            )

        return QueryResponse(
            id=query.id,
            question=query.question,
            answer=query_answers[0] if query_answers else "No Data Found",
            raw_context=raw_contexts[0] if raw_contexts else "No Data Found",
            knowledgeStoreList=knowledge_stores,
        )

    def fetch_llm_response(self, qa_agent, query, knowledge_data):
        self.log.info(f"Fetching LLM response for query: {query.question}")
        return qa_agent.run_query(
            query=query.question,
            collection_name=str(knowledge_data.id),
            document_ids=knowledge_data.documentIds,
            chat_history=query.historyList,
        )

    def is_incorrect_response(self, grader_response):
        return (
            grader_response.validation.lower() == "incorrect"
            or grader_response.answer.lower() == "no data found"
        )

    def get_llm_response_by_processing_whole_document(self, doc_id_list, query, grader):
        doc_id = doc_id_list[0]
        source_file_path = self.document_repo.get_source_path_by_doc_id(doc_id=doc_id)

        file_path = self.s3_helper.download_to(
            source_file_path, os.path.join(SAVE_DIR, os.path.basename(source_file_path))
        )

        response = process_entire_document(
            file_path=file_path,
            query=query,
            source_file_path=source_file_path,
            doc_id=doc_id,
        )

        self.log.info(f"Response from processing entire doc (ID: {doc_id}): {response}")

        grader_response = grader(query.question, response)
        self.log.info(
            f"Grader response for entire doc (ID: {doc_id}): {grader_response}"
        )

        if os.path.exists(file_path):
            os.remove(file_path)

        return response, grader_response, doc_id_list
